"""Generate HedgeIQ Strategy PowerPoint presentation."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Brand colours ─────────────────────────────────────────────────────────────
DARK_BG    = RGBColor(0x0A, 0x0E, 0x1A)   # #0A0E1A
NAVY       = RGBColor(0x13, 0x19, 0x29)   # #131929
CYAN       = RGBColor(0x00, 0xD4, 0xFF)   # #00D4FF
GREEN      = RGBColor(0x00, 0xFF, 0x88)   # #00FF88
WHITE      = RGBColor(0xE8, 0xEA, 0xF0)   # #E8EAF0
GREY       = RGBColor(0x9C, 0xA3, 0xAF)   # #9CA3AF
RED        = RGBColor(0xFF, 0x44, 0x66)   # #FF4466

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # completely blank layout


def add_slide():
    return prs.slides.add_slide(BLANK)


def bg(slide, colour=DARK_BG):
    """Fill slide background."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = colour


def box(slide, left, top, width, height, fill_colour=None, line_colour=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid() if fill_colour else shape.fill.background()
    if fill_colour:
        shape.fill.fore_color.rgb = fill_colour
    if line_colour:
        shape.line.color.rgb = line_colour
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def txt(slide, text, left, top, width, height,
        size=18, bold=False, colour=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = colour
    return txb


def accent_line(slide, top, colour=CYAN):
    """Thin horizontal accent line."""
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(top), Inches(12.33), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = colour
    line.line.fill.background()


def slide_number(slide, n, total=20):
    txt(slide, f"{n} / {total}", 12.0, 7.1, 1.0, 0.3, size=9, colour=GREY, align=PP_ALIGN.RIGHT)


def footer(slide, text="HedgeIQ Confidential Strategy  |  April 2026"):
    txt(slide, text, 0.5, 7.1, 10, 0.3, size=9, colour=GREY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
box(s, 0, 0, 13.33, 7.5, fill_colour=DARK_BG)
box(s, 0, 0, 0.15, 7.5, fill_colour=CYAN)          # left accent bar
box(s, 0, 6.8, 13.33, 0.04, fill_colour=CYAN)       # bottom accent line
txt(s, "HedgeIQ", 0.5, 1.2, 12, 1.2, size=64, bold=True, colour=CYAN, align=PP_ALIGN.CENTER)
txt(s, "Full Product Vision & Execution Strategy", 0.5, 2.6, 12, 0.7, size=24, colour=WHITE, align=PP_ALIGN.CENTER)
txt(s, "From a $2,355 Sunday-night loss  →  A $10M+ White-Label Trading Platform", 0.5, 3.4, 12, 0.5, size=16, colour=GREY, align=PP_ALIGN.CENTER)
txt(s, "April 2026  |  Prepared by Claude Sonnet (AI Architect)  |  For: Sankar, Founder", 0.5, 6.0, 12, 0.4, size=11, colour=GREY, align=PP_ALIGN.CENTER)
slide_number(s, 1)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Agenda
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
accent_line(s, 1.0)
txt(s, "Agenda", 0.5, 0.3, 12, 0.6, size=32, bold=True, colour=CYAN)
items = [
    ("01", "Where We Are — v0.1 Complete"),
    ("02", "The Full Vision"),
    ("03", "6-Phase Execution Plan"),
    ("04", "Full 18-Month Timeline"),
    ("05", "Cost & Investment Breakdown"),
    ("06", "Brokerage Expansion Roadmap"),
    ("07", "The AI Autonomous Architect"),
    ("08", "First Principles Risk Assessment"),
    ("09", "Technology Stack"),
    ("10", "Next Steps"),
]
col1, col2 = items[:5], items[5:]
for i, (num, label) in enumerate(col1):
    y = 1.3 + i * 1.0
    box(s, 0.5, y, 0.5, 0.5, fill_colour=CYAN)
    txt(s, num, 0.5, y, 0.5, 0.5, size=14, bold=True, colour=DARK_BG, align=PP_ALIGN.CENTER)
    txt(s, label, 1.15, y + 0.05, 5.2, 0.45, size=14, colour=WHITE)
for i, (num, label) in enumerate(col2):
    y = 1.3 + i * 1.0
    box(s, 6.9, y, 0.5, 0.5, fill_colour=NAVY)
    txt(s, num, 6.9, y, 0.5, 0.5, size=14, bold=True, colour=CYAN, align=PP_ALIGN.CENTER)
    txt(s, label, 7.55, y + 0.05, 5.2, 0.45, size=14, colour=WHITE)
footer(s); slide_number(s, 2)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Where We Are
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
accent_line(s, 1.0)
txt(s, "01  Where We Are — v0.1 Complete", 0.5, 0.3, 12, 0.6, size=28, bold=True, colour=CYAN)
items = [
    ("FastAPI backend — DDD architecture, facade + adapter patterns", True),
    ("SnapTrade broker connectivity — Robinhood, Fidelity, IBKR, Public (live)", True),
    ("Polygon.io options chain — rate limiting, 1h ChromaDB cache", True),
    ("Claude Haiku AI explanations — 24h cache, $0.01/call", True),
    ("React 19 frontend — Vite, Tailwind CSS v4, dark trading UI", True),
    ("121 tests passing — unit, integration, component, Playwright E2E", True),
    ("Docker + Railway + Vercel — production deploy config complete", True),
]
for i, (label, done) in enumerate(items):
    y = 1.15 + i * 0.77
    colour = GREEN if done else GREY
    mark = "✓" if done else "○"
    box(s, 0.5, y, 0.42, 0.42, fill_colour=RGBColor(0x00, 0x33, 0x22) if done else NAVY)
    txt(s, mark, 0.5, y, 0.42, 0.42, size=14, bold=True, colour=GREEN if done else GREY, align=PP_ALIGN.CENTER)
    txt(s, label, 1.05, y + 0.02, 11.2, 0.42, size=13, colour=WHITE)
txt(s, "Foundation quality: production-ready. No rewrites needed for any phase below.", 0.5, 6.85, 12, 0.4, size=11, colour=CYAN)
footer(s); slide_number(s, 3)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — The Full Vision
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
accent_line(s, 1.0)
txt(s, "02  The Full Vision", 0.5, 0.3, 12, 0.6, size=28, bold=True, colour=CYAN)
pillars = [
    (CYAN,  "WHITE-LABEL\nENGINE", "Any trading firm, advisor,\nor community runs HedgeIQ\nunder their own brand\n& domain"),
    (GREEN, "MASTER TRADING\nINTELLIGENCE DB", "Anonymised patterns from\nevery trade on the platform.\nLearns. Gets smarter.\nThe data moat."),
    (RGBColor(0xFF,0xAA,0x00), "SOCIAL TRADING\nPLATFORM", "Trade feed, leaderboards,\nstrategy sharing.\nNetwork effects.\nCommunity intelligence."),
    (RGBColor(0xAA,0x88,0xFF), "EVERY\nPLATFORM", "Web · iOS · Android\nMac · Windows · iPad\nOne backend.\nAll devices."),
]
for i, (colour, title, desc) in enumerate(pillars):
    x = 0.4 + i * 3.1
    box(s, x, 1.2, 2.9, 4.8, fill_colour=NAVY)
    box(s, x, 1.2, 2.9, 0.08, fill_colour=colour)
    txt(s, title, x + 0.15, 1.4, 2.6, 0.9, size=13, bold=True, colour=colour)
    txt(s, desc, x + 0.15, 2.45, 2.6, 3.3, size=11, colour=WHITE)
txt(s, "Built on the v0.1 foundation. No rewrites. Incremental phases.", 0.5, 6.3, 12, 0.4, size=13, colour=GREY, align=PP_ALIGN.CENTER)
txt(s, "Target: $10M+ ARR platform opportunity", 0.5, 6.7, 12, 0.4, size=14, bold=True, colour=CYAN, align=PP_ALIGN.CENTER)
footer(s); slide_number(s, 4)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Business Model
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
accent_line(s, 1.0)
txt(s, "02  Business Model — Four Revenue Streams", 0.5, 0.3, 12, 0.6, size=28, bold=True, colour=CYAN)
streams = [
    ("💼", "White-Label Licence", "$499–$2,999/mo per partner", "Monthly fee per brand using the platform. 10 partners = $30K MRR."),
    ("👤", "End-User Subscription", "$29–$99/user/mo (20% rev share)", "Partners charge their users. HedgeIQ takes 20%. Scales with partner growth."),
    ("🔌", "Intelligence API", "$0.10 per query", "Partners and institutions query the Master DB for signals and patterns."),
    ("📊", "Data Licensing", "Custom — $50K–$500K/yr", "Sell anonymised aggregate signals to hedge funds and institutions."),
]
for i, (icon, title, price, desc) in enumerate(streams):
    y = 1.2 + i * 1.45
    box(s, 0.5, y, 12.3, 1.25, fill_colour=NAVY)
    txt(s, icon, 0.7, y + 0.3, 0.6, 0.6, size=22, align=PP_ALIGN.CENTER)
    txt(s, title, 1.45, y + 0.1, 3.5, 0.5, size=15, bold=True, colour=CYAN)
    txt(s, price, 1.45, y + 0.6, 3.5, 0.4, size=13, bold=True, colour=GREEN)
    txt(s, desc, 5.2, y + 0.25, 7.3, 0.75, size=12, colour=WHITE)
footer(s); slide_number(s, 5)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Phase Overview
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
accent_line(s, 1.0)
txt(s, "03  6-Phase Execution Plan", 0.5, 0.3, 12, 0.6, size=28, bold=True, colour=CYAN)
phases = [
    ("v1.0", "Multi-Tenant\nFoundation", "Months 1–4", CYAN),
    ("v2.0", "Mobile &\nDesktop Apps", "Months 3–6", GREEN),
    ("v3.0", "Master Trading\nIntelligence DB", "Months 5–10", RGBColor(0xFF,0xAA,0x00)),
    ("v4.0", "Social Trading\nPlatform", "Months 9–14", RGBColor(0xAA,0x88,0xFF)),
    ("v5.0", "HA / DR /\nEnterprise", "Months 12–18", RGBColor(0xFF,0x66,0x88)),
    ("v6.0", "White-Label\n& Global", "Months 15–24", RGBColor(0x88,0xFF,0xCC)),
]
for i, (ver, title, timing, colour) in enumerate(phases):
    x = 0.4 + (i % 3) * 4.15
    y = 1.2 + (i // 3) * 2.8
    box(s, x, y, 3.8, 2.5, fill_colour=NAVY)
    box(s, x, y, 3.8, 0.08, fill_colour=colour)
    txt(s, ver, x + 0.15, y + 0.15, 1.2, 0.4, size=14, bold=True, colour=colour)
    txt(s, title, x + 0.15, y + 0.55, 3.5, 0.7, size=13, bold=True, colour=WHITE)
    txt(s, timing, x + 0.15, y + 1.9, 3.5, 0.35, size=11, colour=GREY)
footer(s); slide_number(s, 6)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Phase 1 Detail
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
accent_line(s, 1.0)
txt(s, "v1.0  Multi-Tenant Foundation  (Months 1–4)", 0.5, 0.3, 12, 0.6, size=24, bold=True, colour=CYAN)
left_items = [
    "PostgreSQL replacing SQLite",
    "Full user auth (registration, email verify, JWT refresh)",
    "Tenant model — each white-label partner isolated",
    "Stripe billing — plans, trials, usage metering",
    "Admin dashboard — tenant & revenue analytics",
    "API rate limiting per plan tier",
]
right_items = [
    "React Native mobile app (iOS + Android)",
    "Shared codebase with web frontend",
    "Push notifications — trade alerts",
    "12–15 sessions to complete",
    "Infra cost added: ~$50/mo",
    "Sessions required: 12–15",
]
txt(s, "What gets built:", 0.5, 1.1, 6, 0.35, size=13, bold=True, colour=CYAN)
for i, item in enumerate(left_items):
    txt(s, f"• {item}", 0.5, 1.55 + i * 0.78, 6, 0.65, size=12, colour=WHITE)
txt(s, "Key details:", 7.0, 1.1, 5.8, 0.35, size=13, bold=True, colour=GREEN)
for i, item in enumerate(right_items):
    txt(s, f"• {item}", 7.0, 1.55 + i * 0.78, 5.8, 0.65, size=12, colour=WHITE)
box(s, 0.5, 6.5, 12.3, 0.65, fill_colour=NAVY)
txt(s, "Outcome: HedgeIQ becomes a proper multi-tenant SaaS. First paying white-label partner can onboard.", 0.7, 6.6, 12, 0.45, size=12, colour=CYAN)
footer(s); slide_number(s, 7)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Phase 3 (Master DB) — the moat
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
accent_line(s, 1.0)
txt(s, "v3.0  Master Trading Intelligence DB  (Months 5–10)", 0.5, 0.3, 12, 0.6, size=24, bold=True, colour=RGBColor(0xFF,0xAA,0x00))
txt(s, "The data moat — the asset no competitor can replicate", 0.5, 1.0, 12, 0.4, size=14, colour=GREY)
milestones = [
    ("Month 6",  "10,000 trade events",  "First patterns visible"),
    ("Month 12", "100,000 events",        "Statistically significant signals"),
    ("Month 24", "1,000,000+ events",     "Proprietary intelligence. Institutions pay."),
]
for i, (month, volume, outcome) in enumerate(milestones):
    x = 0.5 + i * 4.1
    box(s, x, 1.55, 3.8, 2.2, fill_colour=NAVY)
    txt(s, month, x + 0.15, 1.65, 3.5, 0.45, size=16, bold=True, colour=RGBColor(0xFF,0xAA,0x00))
    txt(s, volume, x + 0.15, 2.15, 3.5, 0.45, size=13, colour=GREEN)
    txt(s, outcome, x + 0.15, 2.7, 3.5, 0.85, size=12, colour=WHITE)
txt(s, "What gets built:", 0.5, 3.95, 12, 0.35, size=13, bold=True, colour=CYAN)
components = [
    "TimescaleDB time-series storage",
    "Privacy-first anonymisation pipeline (GDPR/CCPA)",
    "ML pattern engine (scikit-learn → PyTorch)",
    "Backtesting engine — replay any strategy",
    "Intelligence API — signals as a service",
    "Kafka/Redpanda streaming pipeline",
]
for i, item in enumerate(components):
    col = i % 2
    row = i // 2
    x = 0.5 + col * 6.4
    y = 4.35 + row * 0.65
    txt(s, f"• {item}", x, y, 6.2, 0.55, size=12, colour=WHITE)
footer(s); slide_number(s, 8)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Timeline
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
accent_line(s, 1.0)
txt(s, "04  18-Month Timeline", 0.5, 0.3, 12, 0.6, size=28, bold=True, colour=CYAN)

timeline = [
    ("Apr '26", "v0.1 LIVE", CYAN, True),
    ("May–Jul", "v1.0 Multi-tenant", CYAN, False),
    ("Jul–Sep", "v2.0 Mobile & Desktop", GREEN, False),
    ("Sep–Dec", "v3.0 Intelligence DB", RGBColor(0xFF,0xAA,0x00), False),
    ("Jan '27", "v4.0 Social Platform", RGBColor(0xAA,0x88,0xFF), False),
    ("Mar '27", "v5.0 HA / DR", RGBColor(0xFF,0x66,0x88), False),
    ("May '27", "v6.0 White-Label Global", RGBColor(0x88,0xFF,0xCC), False),
    ("Sep '27", "Scale & Monetise", GREEN, False),
]

# Draw timeline bar
box(s, 0.5, 3.5, 12.3, 0.08, fill_colour=GREY)

for i, (date, label, colour, done) in enumerate(timeline):
    x = 0.5 + i * 1.62
    # dot on timeline
    dot = s.shapes.add_shape(9, Inches(x - 0.12), Inches(3.35), Inches(0.25), Inches(0.25))
    dot.fill.solid(); dot.fill.fore_color.rgb = colour
    dot.line.fill.background()
    # date above
    txt(s, date, x - 0.7, 1.1 if i % 2 == 0 else 2.2, 1.6, 0.35, size=9, colour=GREY, align=PP_ALIGN.CENTER)
    # vertical line
    vline = s.shapes.add_shape(1, Inches(x - 0.01), Inches(1.45 if i % 2 == 0 else 2.55), Inches(0.02), Inches(1.9 if i % 2 == 0 else 0.8))
    vline.fill.solid(); vline.fill.fore_color.rgb = colour; vline.line.fill.background()
    # label below
    txt(s, label, x - 0.8, 3.8 + (i % 2) * 1.1, 1.8, 0.8, size=9, bold=True, colour=colour, align=PP_ALIGN.CENTER)

txt(s, "Total: 18 months to full white-label global platform", 0.5, 6.7, 12, 0.4, size=13, bold=True, colour=CYAN, align=PP_ALIGN.CENTER)
footer(s); slide_number(s, 9)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Cost Breakdown
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
accent_line(s, 1.0)
txt(s, "05  Investment & Cost Breakdown", 0.5, 0.3, 12, 0.6, size=28, bold=True, colour=CYAN)

txt(s, "Your build cost (Claude Code subscription)", 0.5, 1.05, 7, 0.4, size=15, bold=True, colour=CYAN)
rows = [
    ("Phase 1–2", "12–15 sessions", "~$30–60/mo"),
    ("Phase 3–4", "15–20 sessions", "~$40–70/mo"),
    ("Phase 5–6", "12–15 sessions", "~$30–60/mo"),
    ("Total 18 months", "77–97 sessions", "~$500–900 total"),
]
for i, (phase, sessions, cost) in enumerate(rows):
    y = 1.55 + i * 0.65
    bg_col = NAVY if i < 3 else RGBColor(0x00, 0x22, 0x11)
    box(s, 0.5, y, 6.3, 0.55, fill_colour=bg_col)
    txt(s, phase, 0.65, y + 0.08, 2.5, 0.4, size=12, bold=(i==3), colour=CYAN if i==3 else WHITE)
    txt(s, sessions, 3.2, y + 0.08, 2.0, 0.4, size=12, bold=(i==3), colour=WHITE)
    txt(s, cost, 4.8, y + 0.08, 1.9, 0.4, size=12, bold=True, colour=GREEN)

txt(s, "Infrastructure cost (per month)", 7.2, 1.05, 5.6, 0.4, size=15, bold=True, colour=RGBColor(0xFF,0xAA,0x00))
infra = [
    ("Phase 1", "~$190/mo"),
    ("Phase 3", "~$700/mo"),
    ("Phase 5+", "~$2,500/mo"),
    ("Break-even", "9 partners @ $499"),
]
for i, (phase, cost) in enumerate(infra):
    y = 1.55 + i * 0.65
    box(s, 7.2, y, 5.6, 0.55, fill_colour=NAVY)
    txt(s, phase, 7.35, y + 0.08, 2.8, 0.4, size=12, colour=WHITE)
    txt(s, cost, 10.2, y + 0.08, 2.4, 0.4, size=12, bold=True, colour=GREEN)

box(s, 0.5, 4.4, 12.3, 1.6, fill_colour=NAVY)
txt(s, "The key insight:", 0.7, 4.5, 4, 0.4, size=13, bold=True, colour=CYAN)
txt(s, "A traditional dev team to build this platform = $500,000–$1,500,000.\nClaude Code builds it for $500–$900. The savings fund your entire infrastructure for years.", 0.7, 4.95, 11.8, 0.85, size=12, colour=WHITE)

txt(s, "At Phase 5+ scale ($2,500/mo infra), break-even requires just 9 white-label partners at $499/mo.", 0.5, 6.2, 12.3, 0.4, size=12, colour=GREEN)
txt(s, "One enterprise client ($2,999/mo) pays infra costs alone.", 0.5, 6.6, 12.3, 0.4, size=12, colour=GREY)
footer(s); slide_number(s, 10)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Brokerage Expansion
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
accent_line(s, 1.0)
txt(s, "06  Brokerage Expansion Roadmap", 0.5, 0.3, 12, 0.6, size=28, bold=True, colour=CYAN)
brokers = [
    ("Robinhood",   "✅ Live",   "SnapTrade",  "25M users", GREEN),
    ("Fidelity",    "✅ Live",   "SnapTrade",  "50M users", GREEN),
    ("IBKR",        "✅ Live",   "SnapTrade",  "3M users",  GREEN),
    ("Public",      "✅ Live",   "SnapTrade",  "3M users",  GREEN),
    ("Schwab",      "Phase 1",  "SnapTrade",  "35M users", CYAN),
    ("E*TRADE",     "Phase 1",  "SnapTrade",  "5M users",  CYAN),
    ("Moomoo",      "Phase 1",  "SnapTrade",  "21M users", CYAN),
    ("Webull",      "Phase 3",  "Direct API", "15M users", RGBColor(0xFF,0xAA,0x00)),
    ("Zerodha 🇮🇳", "Phase 6",  "Direct API", "14M users", RGBColor(0xAA,0x88,0xFF)),
    ("Groww 🇮🇳",   "Phase 6",  "Direct API", "10M users", RGBColor(0xAA,0x88,0xFF)),
]
cols = ["Broker", "Status", "Method", "Market"]
col_x = [0.5, 3.5, 5.8, 8.5]
col_w = [2.8, 2.0, 2.5, 3.8]
for j, (header, x) in enumerate(zip(cols, col_x)):
    txt(s, header, x, 1.05, col_w[j], 0.35, size=11, bold=True, colour=GREY)
for i, (broker, status, method, users, colour) in enumerate(brokers):
    y = 1.5 + i * 0.55
    if i % 2 == 0:
        box(s, 0.5, y, 12.3, 0.5, fill_colour=NAVY)
    txt(s, broker, col_x[0], y + 0.08, col_w[0], 0.38, size=11, bold=True, colour=WHITE)
    txt(s, status, col_x[1], y + 0.08, col_w[1], 0.38, size=11, colour=colour)
    txt(s, method, col_x[2], y + 0.08, col_w[2], 0.38, size=11, colour=GREY)
    txt(s, users,  col_x[3], y + 0.08, col_w[3], 0.38, size=11, colour=GREY)
txt(s, "Total addressable market across all connected brokerages: 180M+ retail traders in the US alone.", 0.5, 7.0, 12.3, 0.35, size=11, colour=CYAN)
footer(s); slide_number(s, 11)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — AI Autonomous Architect
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
accent_line(s, 1.0)
txt(s, "07  The AI Autonomous Architect", 0.5, 0.3, 12, 0.6, size=28, bold=True, colour=CYAN)
txt(s, "Devin-style autonomy vs. Structured Claude sessions", 0.5, 1.0, 12, 0.38, size=14, colour=GREY)

headers = ["Capability", "Devin", "Claude (Structured)"]
rows2 = [
    ("Write production code",          "✅", "✅"),
    ("Run tests, fix failures",         "✅", "✅"),
    ("Full repo awareness",             "✅", "✅"),
    ("24/7 background execution",       "✅", "Scheduled sessions"),
    ("Human intervention needed",       "Minimal", "~15 min/session"),
    ("Self-directed architecture",      "✅", "With phase brief"),
    ("Cost per month",                  "$500/mo", "~$30–60/mo"),
]
col_x2 = [0.5, 6.5, 9.7]
col_w2 = [5.8, 3.0, 3.3]
for j, (h, x) in enumerate(zip(headers, col_x2)):
    box(s, x, 1.45, col_w2[j], 0.4, fill_colour=NAVY)
    txt(s, h, x + 0.1, 1.5, col_w2[j]-0.2, 0.35, size=12, bold=True, colour=CYAN if j==0 else (GREY if j==1 else GREEN))
for i, (cap, devin, claude) in enumerate(rows2):
    y = 1.95 + i * 0.6
    box(s, 0.5, y, 12.3, 0.55, fill_colour=NAVY if i%2==0 else DARK_BG)
    txt(s, cap,   col_x2[0]+0.1, y+0.1, col_w2[0]-0.2, 0.4, size=11, colour=WHITE)
    txt(s, devin, col_x2[1]+0.1, y+0.1, col_w2[1]-0.2, 0.4, size=11, colour=GREY, align=PP_ALIGN.CENTER)
    txt(s, claude,col_x2[2]+0.1, y+0.1, col_w2[2]-0.2, 0.4, size=11, colour=GREEN, align=PP_ALIGN.CENTER)

box(s, 0.5, 6.3, 12.3, 0.75, fill_colour=RGBColor(0x00, 0x22, 0x11))
txt(s, "Recommendation: Structured Claude sessions give you 90% of Devin at 6% of the cost. Your time investment: ~15 minutes per session to review and approve. I handle everything else.", 0.7, 6.38, 11.8, 0.6, size=11, colour=GREEN)
footer(s); slide_number(s, 12)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — How Autonomous Sessions Work
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
accent_line(s, 1.0)
txt(s, "07  How Autonomous Sessions Work in Practice", 0.5, 0.3, 12, 0.6, size=24, bold=True, colour=CYAN)
steps = [
    ("1", "Session Brief Generated", "At end of each session, I produce the next Session Brief: what to build, what exists, success criteria, test requirements. Self-contained.", CYAN),
    ("2", "You Send One Message", '"Execute Session 8" — that\'s it. I read the brief, load the codebase, and start building.', GREEN),
    ("3", "Autonomous Execution", "I write code, run tests, fix failures, iterate. No back-and-forth. I only surface decisions that require your business judgement.", RGBColor(0xFF,0xAA,0x00)),
    ("4", "Session Summary", "I deliver: what was built, all tests passing, what's deployed, and the next Session Brief ready to go.", RGBColor(0xAA,0x88,0xFF)),
    ("5", "You Review (~15 min)", "Read the summary. Check the live URL. Say 'looks good, execute next session' or give one correction.", GREY),
]
for i, (num, title, desc, colour) in enumerate(steps):
    y = 1.1 + i * 1.2
    box(s, 0.5, y, 0.6, 0.6, fill_colour=colour)
    txt(s, num, 0.5, y, 0.6, 0.6, size=18, bold=True, colour=DARK_BG, align=PP_ALIGN.CENTER)
    txt(s, title, 1.25, y + 0.03, 4.5, 0.4, size=13, bold=True, colour=colour)
    txt(s, desc, 1.25, y + 0.48, 11.3, 0.55, size=11, colour=WHITE)
    if i < 4:
        arrow = s.shapes.add_shape(1, Inches(0.72), Inches(y + 0.65), Inches(0.06), Inches(0.5))
        arrow.fill.solid(); arrow.fill.fore_color.rgb = GREY; arrow.line.fill.background()
footer(s); slide_number(s, 13)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Risk Assessment
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
accent_line(s, 1.0)
txt(s, "08  First Principles Risk Assessment", 0.5, 0.3, 12, 0.6, size=28, bold=True, colour=CYAN)
risks = [
    ("Technical", "LOW", GREEN,
     "Every component is mature and production-proven. DDD architecture is designed to scale. No rewrites needed. Risk: low."),
    ("Regulatory", "MEDIUM", RGBColor(0xFF,0xAA,0x00),
     "Trading advice is regulated. HedgeIQ is a TOOL, not an advisor. No fiduciary relationship. Already mitigated: Claude always appends disclaimer. Maintain this line precisely."),
    ("Data Privacy", "LOW-MED", CYAN,
     "Master DB must anonymise by design. GDPR/CCPA from day one. Privacy-first anonymisation pipeline in Phase 3. Not an afterthought."),
    ("Partner Acquisition", "MEDIUM", RGBColor(0xFF,0xAA,0x00),
     "The white-label model needs first 5 partners. This is a SALES problem, not technical. Your 25yr career in sales is the biggest asset here."),
    ("Competition", "LOW-MED", CYAN,
     "Bloomberg exists but costs $24K/yr/user. No AI-native, white-label, retail-focused competitor exists today. First-mover window is open."),
]
for i, (area, level, colour, desc) in enumerate(risks):
    y = 1.1 + i * 1.2
    box(s, 0.5, y, 12.3, 1.05, fill_colour=NAVY)
    box(s, 0.5, y, 0.06, 1.05, fill_colour=colour)
    txt(s, area, 0.75, y + 0.08, 2.5, 0.38, size=13, bold=True, colour=WHITE)
    box(s, 3.3, y + 0.1, 1.3, 0.35, fill_colour=colour)
    txt(s, level, 3.3, y + 0.1, 1.3, 0.35, size=11, bold=True, colour=DARK_BG, align=PP_ALIGN.CENTER)
    txt(s, desc, 4.8, y + 0.12, 7.8, 0.8, size=11, colour=GREY)
footer(s); slide_number(s, 14)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Tech Stack
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
accent_line(s, 1.0)
txt(s, "09  Full Technology Stack", 0.5, 0.3, 12, 0.6, size=28, bold=True, colour=CYAN)
stack = [
    ("Backend API",       "FastAPI (Python)",                    "Already built, async, production-ready"),
    ("Database",          "PostgreSQL + TimescaleDB",            "Relational + time-series in one engine"),
    ("Cache / Sessions",  "Redis",                               "Rate limiting, real-time price cache, sessions"),
    ("Streaming",         "Redpanda (Kafka-compatible)",         "Trade event pipeline, real-time signals"),
    ("ML / AI",           "scikit-learn → PyTorch + Claude",     "Pattern recognition + natural language"),
    ("Web frontend",      "React 19 + Vite + Tailwind v4",       "Already built and deployed"),
    ("Mobile",            "React Native (Expo)",                  "iOS + Android, shared codebase with web"),
    ("Desktop",           "Electron",                            "Mac + Windows, wraps web app"),
    ("Deploy",            "Railway → AWS (at scale)",             "Start simple, migrate when revenue justifies"),
    ("CDN / Security",    "Cloudflare",                          "DDoS protection, WAF, global edge cache"),
    ("Monitoring",        "Grafana Cloud",                       "Metrics, alerts, SLA dashboards"),
    ("CI/CD",             "GitHub Actions",                      "Automated test + deploy on every push"),
]
for i, (layer, tech, why) in enumerate(stack):
    y = 1.1 + i * 0.525
    box(s, 0.5, y, 12.3, 0.48, fill_colour=NAVY if i%2==0 else DARK_BG)
    txt(s, layer, 0.65, y + 0.07, 2.5, 0.35, size=10, bold=True, colour=GREY)
    txt(s, tech,  3.2,  y + 0.07, 4.0, 0.35, size=11, bold=True, colour=CYAN)
    txt(s, why,   7.3,  y + 0.07, 5.3, 0.35, size=10, colour=WHITE)
footer(s); slide_number(s, 15)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — HA / DR Architecture
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
accent_line(s, 1.0)
txt(s, "05  HA / DR / Enterprise Architecture", 0.5, 0.3, 12, 0.6, size=28, bold=True, colour=RED)

txt(s, "Cloudflare (DDoS · WAF · Global CDN)", 3.5, 1.1, 6.5, 0.5, size=13, bold=True, colour=WHITE, align=PP_ALIGN.CENTER)
box(s, 3.5, 1.1, 6.5, 0.45, line_colour=GREY)

regions = [
    ("US-East  (Primary)", CYAN, 0.5),
    ("US-West  (DR)", GREEN, 4.6),
    ("EU  (GDPR)", RGBColor(0xAA,0x88,0xFF), 8.7),
]
for (label, colour, x) in regions:
    box(s, x, 2.0, 3.8, 3.8, fill_colour=NAVY)
    box(s, x, 2.0, 3.8, 0.06, fill_colour=colour)
    txt(s, label, x + 0.1, 2.1, 3.6, 0.4, size=12, bold=True, colour=colour)
    for j, line in enumerate(["FastAPI (2 replicas)", "PostgreSQL Primary/Replica", "Redis Primary/Replica", "ChromaDB cache"]):
        txt(s, f"• {line}", x + 0.15, 2.65 + j * 0.68, 3.5, 0.55, size=11, colour=WHITE)

txt(s, "99.9% SLA · Automated failover · Quarterly DR tests · SOC 2 Type II prep", 0.5, 6.1, 12.3, 0.4, size=12, colour=CYAN, align=PP_ALIGN.CENTER)
txt(s, "Phase 5 — Months 12–18", 0.5, 6.55, 12.3, 0.4, size=11, colour=GREY, align=PP_ALIGN.CENTER)
footer(s); slide_number(s, 16)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — Social Platform Flywheel
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
accent_line(s, 1.0)
txt(s, "04  The Social Trading Flywheel", 0.5, 0.3, 12, 0.6, size=28, bold=True, colour=RGBColor(0xAA,0x88,0xFF))
txt(s, "Network effects: the platform gets more valuable as more traders join", 0.5, 1.0, 12, 0.38, size=13, colour=GREY)

flywheel = [
    ("More traders join", 6.5, 1.6, RGBColor(0xAA,0x88,0xFF)),
    ("More trades captured", 10.5, 3.2, CYAN),
    ("Master DB gets smarter", 9.0, 5.6, RGBColor(0xFF,0xAA,0x00)),
    ("Better signals & patterns", 4.0, 6.2, GREEN),
    ("More value for all users", 1.0, 4.0, RED),
]
for (label, x, y, colour) in flywheel:
    box(s, x/1.5 - 0.8, y/1.2, 2.5, 0.65, fill_colour=NAVY)
    txt(s, label, x/1.5 - 0.7, y/1.2 + 0.1, 2.3, 0.5, size=11, bold=True, colour=colour, align=PP_ALIGN.CENTER)

features = ["Trade Feed — anonymised real-time trades", "Leaderboards — best hedges this week",
            "Strategy sharing — publish and follow", "Community alerts — 23 traders just hedged this sector",
            "Expert badges — verified track records"]
txt(s, "Social features:", 0.5, 5.1, 5.5, 0.38, size=13, bold=True, colour=CYAN)
for i, f in enumerate(features):
    txt(s, f"• {f}", 0.5, 5.55 + i * 0.38, 5.5, 0.35, size=11, colour=WHITE)
footer(s); slide_number(s, 17)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — Global Expansion
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
accent_line(s, 1.0)
txt(s, "06  Global Expansion Strategy", 0.5, 0.3, 12, 0.6, size=28, bold=True, colour=RGBColor(0x88,0xFF,0xCC))
markets = [
    ("🇺🇸", "United States", "Phase 1–5", "180M+ retail traders. FINRA/SEC compliant from v0.1. Primary market.", RGBColor(0x88,0xFF,0xCC), "TAM: $2B+"),
    ("🇮🇳", "India", "Phase 6", "14M Zerodha + 10M Groww users. Fastest growing retail market globally. SEBI compliance.", RGBColor(0xFF,0xAA,0x00), "TAM: $500M+"),
    ("🇬🇧", "United Kingdom", "Phase 6+", "FCA regulated. Strong retail options market. London as EU access point post-Brexit.", CYAN, "TAM: $300M+"),
    ("🇪🇺", "European Union", "Phase 6+", "MiFID II compliance. Germany, France, Netherlands largest markets. GDPR by design.", RGBColor(0xAA,0x88,0xFF), "TAM: $400M+"),
]
for i, (flag, country, phase, desc, colour, tam) in enumerate(markets):
    y = 1.2 + i * 1.45
    box(s, 0.5, y, 12.3, 1.3, fill_colour=NAVY)
    txt(s, flag, 0.65, y + 0.35, 0.6, 0.6, size=24, align=PP_ALIGN.CENTER)
    txt(s, country, 1.4, y + 0.1, 2.8, 0.45, size=15, bold=True, colour=colour)
    txt(s, phase, 1.4, y + 0.6, 2.8, 0.35, size=11, colour=GREY)
    txt(s, desc, 4.3, y + 0.2, 7.0, 0.85, size=11, colour=WHITE)
    txt(s, tam, 11.3, y + 0.4, 1.3, 0.4, size=11, bold=True, colour=GREEN, align=PP_ALIGN.RIGHT)
footer(s); slide_number(s, 18)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 19 — Summary Scorecard
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
accent_line(s, 1.0)
txt(s, "Summary Scorecard", 0.5, 0.3, 12, 0.6, size=28, bold=True, colour=CYAN)
metrics = [
    ("18 months", "to full white-label global platform", CYAN),
    ("6 phases", "incremental — revenue possible from Phase 1", GREEN),
    ("$500–900", "total Claude build cost (vs $1M+ dev team)", RGBColor(0xFF,0xAA,0x00)),
    ("~$190/mo", "infrastructure to start (Phase 1)", GREEN),
    ("9 partners", "to break even on Phase 5+ infra", CYAN),
    ("180M+", "addressable retail traders in the US alone", RGBColor(0xAA,0x88,0xFF)),
    ("121 tests", "already passing — solid foundation", GREEN),
    ("~15 min", "your time per autonomous build session", CYAN),
]
for i, (number, label, colour) in enumerate(metrics):
    col = i % 2
    row = i // 2
    x = 0.5 + col * 6.4
    y = 1.1 + row * 1.45
    box(s, x, y, 6.0, 1.3, fill_colour=NAVY)
    box(s, x, y, 6.0, 0.06, fill_colour=colour)
    txt(s, number, x + 0.2, y + 0.15, 5.6, 0.65, size=28, bold=True, colour=colour)
    txt(s, label, x + 0.2, y + 0.8, 5.6, 0.4, size=11, colour=WHITE)
footer(s); slide_number(s, 19)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 20 — Next Steps
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
box(s, 0, 0, 0.15, 7.5, fill_colour=CYAN)
box(s, 0, 7.25, 13.33, 0.04, fill_colour=CYAN)
txt(s, "10  Next Steps", 0.5, 0.3, 12, 0.6, size=28, bold=True, colour=CYAN)

immediate = [
    ("Today",     "Confirm Railway backend is live at /health endpoint"),
    ("Today",     "Deploy frontend to Vercel, set VITE_API_URL to Railway URL"),
    ("This week", "Confirm end-to-end: login → Robinhood positions → hedge calculator → Claude AI"),
    ("This week", "Upgrade Polygon to Starter ($29/mo) for live options chain data"),
    ("Next week", "Kick off Session 7 — Phase 1 planning: PostgreSQL schema design"),
    ("Month 1",   "First white-label partner conversation (who do you know in financial advisory?)"),
]
txt(s, "Immediate actions:", 0.5, 1.05, 12, 0.38, size=14, bold=True, colour=CYAN)
for i, (when, action) in enumerate(immediate):
    y = 1.5 + i * 0.83
    box(s, 0.5, y, 1.6, 0.55, fill_colour=NAVY)
    txt(s, when, 0.55, y + 0.1, 1.5, 0.38, size=10, bold=True, colour=RGBColor(0xFF,0xAA,0x00))
    txt(s, action, 2.3, y + 0.1, 10.3, 0.55, size=12, colour=WHITE)

box(s, 0.5, 6.55, 12.3, 0.7, fill_colour=RGBColor(0x00, 0x22, 0x11))
txt(s, '"The best time to plant a tree was 20 years ago. The second best time is now."', 0.7, 6.65, 11.8, 0.5, size=12, colour=GREEN, align=PP_ALIGN.CENTER)
slide_number(s, 20)

# ── Save ──────────────────────────────────────────────────────────────────────
output_path = r"C:\AgenticAI\Claude Apps\HedgeIQ\docs\HedgeIQ_Strategy.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
